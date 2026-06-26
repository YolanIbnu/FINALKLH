"""
SIPETA TPK - Seminar Proposal PPT Generator (Part 2: Builder)
Run: python generate_ppt_part2.py
Output: Seminar_Proposal_SIPETA_TPK.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, sys

# Import slide data
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from generate_ppt_part1 import SLIDES

# Color scheme - Professional dark blue theme
BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)      # Dark navy
BG_CARD = RGBColor(0x1B, 0x2A, 0x41)      # Card background
ACCENT = RGBColor(0x00, 0x96, 0xC7)        # Bright blue accent
ACCENT2 = RGBColor(0x48, 0xCA, 0xE4)       # Light blue
GOLD = RGBColor(0xFF, 0xD6, 0x00)          # Gold accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
SUBTITLE_CLR = RGBColor(0x90, 0xE0, 0xEF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=WHITE, line_spacing=1.3, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.space_after = Pt(2)

        # Color logic
        if line.startswith("■") or line.startswith("❶") or line.startswith("❷") or \
           line.startswith("❸") or line.startswith("❹"):
            p.font.color.rgb = ACCENT2
            p.font.bold = True
        elif line.startswith("✓") or line.startswith("→"):
            p.font.color.rgb = GOLD
            p.font.bold = True
        elif any(line.startswith(prefix) for prefix in
                 ["IDENTIFIKASI", "RUMUSAN", "BATASAN", "TUJUAN", "FRONTEND",
                  "PEMETAAN", "BACKEND", "ADMIN", "STAFF", "MANFAAT", "MEKANISME",
                  "RESEARCH GAP", "Alur", "Grade SUS", "Rumus", "No |"]):
            p.font.color.rgb = ACCENT2
            p.font.bold = True
        elif line.startswith("RQ") or line.startswith("T1") or line.startswith("T2") or \
             line.startswith("T3") or line.startswith("T4"):
            p.font.color.rgb = SUBTITLE_CLR
        elif line.startswith("    "):
            p.font.color.rgb = LIGHT_GRAY
            p.font.size = Pt(font_size - 1)
        elif line.startswith("•"):
            p.font.color.rgb = WHITE
        elif line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or \
             line.startswith("4.") or line.startswith("5.") or line.startswith("6."):
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif line == "":
            p.font.size = Pt(6)
        else:
            p.font.color.rgb = color

    return txBox

def build_cover_slide(prs, slide_data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    # Top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), ACCENT)

    # Left accent strip
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT)

    # Side decorative box
    add_accent_bar(slide, Inches(0.3), Inches(1.5), Inches(0.06), Inches(4.5), GOLD)

    if slide_num == 1:
        # Title
        add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.6),
                    slide_data["title"], 28, ACCENT2, True)

        # Main title - split by newlines
        lines = slide_data["subtitle"].split("\n")
        y = 1.6
        for line in lines:
            add_textbox(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.6),
                        line, 26, WHITE, True)
            y += 0.55

        # Content lines
        add_multiline_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(3.5),
                              slide_data["content"], 18, LIGHT_GRAY)
    else:
        # Closing slide
        add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                    slide_data["title"], 40, GOLD, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
                    slide_data["subtitle"], 22, ACCENT2, False, PP_ALIGN.CENTER)
        add_multiline_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(4),
                              slide_data["content"], 20, LIGHT_GRAY)

def build_content_slide(prs, slide_data, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)

    # Top bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

    # Title bar background
    add_accent_bar(slide, Inches(0), Inches(0.06), SLIDE_W, Inches(0.9), BG_CARD)

    # Title text
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10.5), Inches(0.7),
                slide_data["title"], 24, WHITE, True)

    # Slide number
    add_textbox(slide, Inches(11.5), Inches(0.2), Inches(1.5), Inches(0.5),
                f"{slide_num}", 16, ACCENT2, True, PP_ALIGN.RIGHT)

    # Subtitle if exists
    content_top = 1.15
    if slide_data.get("subtitle"):
        add_textbox(slide, Inches(0.6), Inches(1.05), Inches(11), Inches(0.4),
                    slide_data["subtitle"], 16, SUBTITLE_CLR, False)
        content_top = 1.5

    # Accent line under title
    add_accent_bar(slide, Inches(0.6), Inches(content_top - 0.05), Inches(2), Inches(0.04), GOLD)

    # Content
    add_multiline_textbox(slide, Inches(0.6), Inches(content_top + 0.1),
                          Inches(12), Inches(6.0 - content_top),
                          slide_data["content"], 16, WHITE)

    # Bottom bar
    add_accent_bar(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), BG_CARD)
    add_textbox(slide, Inches(0.5), Inches(7.35), Inches(8), Inches(0.15),
                "SIPETA TPK — Seminar Proposal Skripsi — Yolan Ibnu Prasetya",
                9, LIGHT_GRAY)

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(SLIDES):
        slide_num = i + 1
        if slide_data["layout"] == "cover":
            build_cover_slide(prs, slide_data, slide_num)
        else:
            build_content_slide(prs, slide_data, slide_num)
        print(f"  ✓ Slide {slide_num}: {slide_data['title']}")

    output_dir = r"D:\idm download\SEMESTER 8 YOLAN\TA YOLAN\SKRIPSI ASLI YOLAN"
    output_path = os.path.join(output_dir, "Seminar_Proposal_SIPETA_TPK.pptx")
    prs.save(output_path)
    print(f"\n✅ PPT berhasil disimpan di:\n   {output_path}")
    print(f"   Total: {len(SLIDES)} slides")

if __name__ == "__main__":
    print("🔨 Generating Seminar Proposal PPT...")
    main()
